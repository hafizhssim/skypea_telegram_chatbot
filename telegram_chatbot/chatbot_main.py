# Standard Libraries
import os
import re
import json
import time
import random
import datetime
import logging
import unicodedata
from random import choice
from itertools import chain
from urllib.parse import urlparse, parse_qs
from concurrent.futures import ThreadPoolExecutor, as_completed
from io import BytesIO
import xml.etree.ElementTree as ETmport 
import psycopg2
import numpy as np
import pandas as pd
import spacy
import openai
import tiktoken
import tensorflow_hub as hub
from bs4 import BeautifulSoup
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.corpus import wordnet
from sklearn.metrics.pairwise import cosine_similarity
from sklearn.neighbors import NearestNeighbors
from pytube import YouTube
from urllib.parse import urlparse
from selenium import webdriver
from termcolor import colored
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader
from youtube_transcript_api import YouTubeTranscriptApi
# from datetime import datetime
import requests
# NLTK Downloads
import nltk
nltk.download('wordnet')
nltk.download('punkt')

# Database Configuration
DB_NAME = os.getenv('DB_NAME')
if DB_NAME is None:
    logging.error("Environment variable DB_NAME is not set.")
DB_USER = os.getenv('DB_USER')
DB_PASSWORD = os.getenv('DB_PASS')
DB_HOST = 'localhost'
DB_PORT = '5432'

# OpenAI Configuration
openai.api_key = os.getenv('OPENAI_API_KEY')

#5 Function to clear the terminal
def clear_terminal():
    # Clear terminal command for Windows
    if os.name == 'nt':
        os.system('cls')

# Clear the terminal before running the code
clear_terminal()

class gpt:

# SETUP

    def __init__(self):
        # Fetch database credentials from OS environment
        db_user = os.getenv('DB_USER')
        db_pass = os.getenv('DB_PASS')

        # Set up database connection
        try:
            with psycopg2.connect(
                dbname="skypea",
                user=db_user,
                password=db_pass,
                host="localhost",
                port="5432"
            ) as self.conn:
                self.cur = self.conn.cursor()
        except psycopg2.Error as e:
            print(colored("Error: Could not make connection to the Postgres database", "red"))
            print(colored({e},"red"))
        
        # Initialize conversation history
        self.conversation_id = None

        # Initialize the current timestamp
        self.current_timestamp = datetime.datetime.now()
        
        # Define the path to the extracted model folder
        model_path = "G:\\My Drive\\skypea\\universal-sentence-encoder_4"
        self.model = hub.load(model_path)   # "https://tfhub.dev/google/universal-sentence-encoder/4"
        self.nlp = spacy.load("en_core_web_sm")

# 2 GENERATE CURRENT QUERY EMBBEDDINGS AND INITIALIZING K-NN MODEL 

    # 2a GENERATE CURRENT QUERY EMBBEDDINGS
    def generate_current_query_embedding(self, query):
        try:
            # Generate embeddings for the user_query
            current_user_query_embedding = self.model([query])[0].numpy()

            return current_user_query_embedding

        except Exception as e:
            print(colored(f"Error: Unable to generate embeddings. {e}", "red"))
            return None

    # 2b INITIALIZING K-NN MODEL
    def initialize_knn_model(self):
        self.cur.execute("""
                            SELECT id, combined_response_embedding
                            FROM skypea_db;
                        """)
        all_embedding_data = self.cur.fetchall()

        # Create a list of embeddings and their corresponding IDs
        all_embeddings = [np.array(combined_response_emb) for id, combined_response_emb in all_embedding_data if combined_response_emb is not None]

        # Determine the number of samples and set k dynamically
        n_samples = len(all_embeddings)  # Changed from len(all_embedding_data) to len(all_embeddings) to only count non-None embeddings
        k = min(10, n_samples)

        # Train k-NN model
        if k > 0:  # Check to make sure there are some samples
            self.knn_model = NearestNeighbors(n_neighbors=k, algorithm='auto').fit(all_embeddings)
            # print(k)
        else:
            print(colored("Not enough samples to fit k-NN model.","red"))

# 3a DATABASE PIPELINE

    def fetch_and_clean_from_db(self, current_user_query_embedding, query):
        # Step 1: Fetch all combined_response data points from the database
        self.cur.execute("""
            SELECT id, combined_response_embedding
            FROM skypea_db WHERE combined_response_embedding IS NOT NULL;
        """)
        all_data = self.cur.fetchall()

        knn_data = []
        for row in all_data:
            id = row[0]
            for i, column in enumerate(['combined_response_embedding']):
                if row[i+1] is not None:
                    knn_data.append((id, column, np.array(row[i+1])))

        unique_data = []
        for id, column, embedding in knn_data:
            is_unique = True
            for _, _, unique_embedding in unique_data:
                similarity = cosine_similarity([embedding], [unique_embedding])
                if similarity > 0.85:  # Adjust this threshold as needed
                    is_unique = False
                    break
            if is_unique:
                unique_data.append((id, column, embedding))

        if len(unique_data) == 0:
            print(colored("No data fetched from database.","red"))
            return [("No Data", "No data fetched from database.")]
        
        n_neighbors = min(5, len(unique_data))

        knn_model = NearestNeighbors(n_neighbors=n_neighbors, algorithm='auto')
        knn_model.fit([x[2] for x in unique_data])
        distances, indices = knn_model.kneighbors([current_user_query_embedding])
        sorted_data = [unique_data[i] for i in indices[0]]

        top_text_data = []
        for id, column, _ in sorted_data:
            text_column = column.replace("_embedding", "")
            self.cur.execute(f"SELECT {text_column} FROM skypea_db WHERE id = %s", (id,))
            text_data = self.cur.fetchone()[0]
            formatted_data = f"id('{id}') - '{text_data}'"
            top_text_data.append(formatted_data)

        concatenated_text_data = ' '.join(top_text_data)
        # print(f'concatenated_text_data = {concatenated_text_data}')

        # Fetch 5 latest entries from the 'summary' column in 'skypea_db'
        self.cur.execute("""
                        SELECT summary FROM skypea_db ORDER BY id DESC LIMIT 10;
                        """)
        previous_conversation = self.cur.fetchall()
        previous_conversation.reverse()
        previous_conversation = str(previous_conversation)
        # print(colored(previous_conversation,"magenta"))

        # OPENAI API CALL
        max_retries = 5
        delay = 1
        for attempt in range(max_retries):
            try:
                response = openai.ChatCompletion.create(
                    model="gpt-3.5-turbo-16k",  # or gpt-4 when available
                    temperature=0.5,
                    messages=[
                        {
                            "role": "system",
                            "content": f"""
                            You are a chatbot with memory access. 
                            You should recall and analyze data from the past conversations as well as the most relevant information from your memory as provided below. 
                            Your role is to assist on the current query and generate in into one components only out of the two below.
                            If theres enough information about it based on your memory, provide the findings in the "memory :" if not, then provide searh query in the "internet_query" to find the asnwer. 

                            One component should have response, one component will have '-' as its output. 
                            
                            Example:
                            
                            Condition where the question has a internet links or asking for the content of the internet links, it must return the internet links in the internet_query output and not memory
                            memory : -
                            internet_query: internet links (any https or .com provided)

                            Condition where there is answer to the question in the database regarding the query,
                            memory : response
                            internet_query: -

                            condition where there no asnwer in the database regarding the query, or the user providing the internet links
                            memory : -
                            internet_query: internet search query that we are going to use.
                            
                            - Think first then decide. most of the time the question always related to the previous conversations, some time it does not. 
                            - Be decisive, dont say im sorry etc. if there no answer to produce for memory, just return interenet_query to provide interenet search.
                            - Limited to only 3000 tokens of output
                            """
                        },
                        {
                            "role": "user",
                            "content": f"""
                            From,
                            (Top 5 Relevant Data from your memory: {concatenated_text_data}.)
                            (Previous conversations you recalled from your memory, always relate the question if the question has no topic or subject that is not provided in the question: {previous_conversation}.)
                            (Now the question is? {query}.)
                            """
                        }
                    ]
                )
                db_data = response.choices[0].message['content'].strip()
                
                # print(colored(db_data,"magenta"))

                if "internet_query:" in db_data and "memory: " in db_data:
                    internet_query = re.search(r"internet_query: (.+)", db_data).group(1)
                    internet_query = str(internet_query)
                    print(colored(f"{internet_query}", "green"))

                    memory = re.search(r"memory: (.+)", db_data).group(1)
                    memory = str(memory)
                    print(colored(f"Memory: {memory}", "green"))

                    if internet_query == '-' or memory == '-':
                        return memory, internet_query
                    else:
                        raise ValueError("If both components are present, one must be '-'.")

                elif "internet_query:" in db_data or "memory: " in db_data:
                    internet_query = None if "internet_query:" not in db_data else re.search(r"internet_query: (.+)", db_data).group(1)
                    memory = None if "memory: " not in db_data else re.search(r"memory: (.+)", db_data).group(1)

                    return memory, internet_query

                else:
                    raise ValueError("Either memory or internet_query should be present.")

            except Exception as e:
                print(colored(f"An error occurred at gpt_query: {e}", "red"))
                if attempt < max_retries - 1:
                    print(colored(f"Retrying... ({attempt + 1})", "red"))
                    time.sleep(delay)
                else:
                    print(colored("Max retries reached. Could not complete the operation.", "red"))
                    return None, None
    
# 3b INTERNET PIPELINE

    def internet_route(self, current_user_query_embedding, query, internet_query):
        if internet_query == '-' or internet_query is None:
            return '-'        
        
        failed_urls = set()
        url_content_list = []

        # Extract URLs from query
        internet_links = re.findall(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', query)

        if internet_links:
            # If there are URLs in the query, process them concurrently
            with ThreadPoolExecutor() as executor:
                # print(internet_links)
                futures = {}
                for link in internet_links:
                    futures[executor.submit(self.fetch_content_and_generate_embedding, link)] = link

                for future in as_completed(futures):
                    link = futures[future]
                    try:
                        cleaned_content, embedding = future.result()
                        url_content_list.append((link, cleaned_content, embedding))
                    except Exception as e:
                        print(f"Exception occurred: {e}")
                        failed_urls.add(link)
        else:
            # If no URLs in the query, perform Google search and fetch top 6 links
            links = self.fetch_google_links(internet_query, 0)
            top_links = links[:5]  # Limit to top n links
            # print(top_links)

            with ThreadPoolExecutor() as executor:
                futures = {}
                for link in top_links:
                    clean_url = parse_qs(urlparse(link).query).get('q', [None])[0]
                    if clean_url and clean_url not in failed_urls:
                        futures[executor.submit(self.fetch_content_and_generate_embedding, clean_url)] = clean_url

                for future in as_completed(futures):
                    clean_url = futures[future]
                    try:
                        cleaned_content, embedding = future.result()
                        url_content_list.append((clean_url, cleaned_content, embedding))
                    except Exception as e:
                        print(f"Exception occurred: {e}")
                        failed_urls.add(clean_url)

        sorted_links = self.sorted_url_content_embedding_list(url_content_list, current_user_query_embedding)
        # print(sorted_links)

        # Step 4: Analyze the top_text_data using OpenAI API
        concatenated_text_data = ' '.join(map(str, sorted_links))

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-16k",  # or gpt-4 when available
            temperature=0.7,
            messages=[
                {
                    "role": "system",
                    "content": f"""
                    You are a chatbot with real-time internet access. If the Internet Query stated 'not relevant', just return 'not relevant'. If there's link/s in it, details its content, by reviewing the internet search results.
                    Your role is to assist in processing the real-time / internet data to keep its contents intact in details as possible while limiting it to only 4000 token.
                    """
                },
                {
                    "role": "user",
                    "content": f"""
                    (Internet Query: {internet_query}.)
                    (Top 5 Results real-time data/ internet search: {concatenated_text_data}.)
                    (Current Query: {query}.)
                    
                    Output format example:
                    (url_1 - title_1 - summarized and detailed content_1
                    url_2 - title_2 - summarized and detailed content_2
                    summary)
                    """
                }
            ]
        )
        sorted_links = response.choices[0].message['content'].strip()
        # print(sorted_links)

        return sorted_links

    def fetch_google_links(self, internet_query, page):
        start_num = page * 10
        google_url = f"https://www.google.com/search?q={internet_query}&start={start_num}"
        session = requests.Session()
        session.headers.update({
            'User-Agent': 'Mozilla/5.0',
            'Referer': 'https://www.google.com/'
        })
        response = session.get(google_url)
        soup = BeautifulSoup(response.text, 'html.parser')
        link_elements = soup.select('.egMi0.kCrYT a')
        links = [link['href'] for link in link_elements]
        return links

    # 3b 2 Fetch and score the internet links
    def fetch_content_and_generate_embedding(self, url):
        text = ""
        max_retries = 3
        delay = 1
        blocked_urls = set()

        user_agents = [
            'Mozilla/5.0 (Windows NT 10.0; Win64; x64)',
            'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_4)',
            'Mozilla/5.0 (X11; Ubuntu; Linux x86_64; rv:62.0) Gecko/20100101 Firefox/62.0',
            'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537',
            'Mozilla/5.0 (iPhone; CPU iPhone OS 11_0 like Mac OS X) AppleWebKit/604.1.38 (KHTML, like Gecko) Version/11.0 Mobile/15A372 Safari/604.1'
        ]

        for attempt in range(max_retries):
            try:
                if url in blocked_urls:
                    print(colored(f"Skipping blocked URL: {url}", "red"))
                    return 0, ''

                headers = {'User-Agent': random.choice(user_agents)}
                response = requests.get(url, headers=headers, timeout=10)

                if response.status_code == 429:  # Handling rate limiting
                    print(colored(f"Rate limit exceeded for URL: {url}. Waiting...", "yellow"))
                    time.sleep(5)  # Wait for 60 seconds before retrying
                    continue  # Skip the rest of the loop and retry

                if response.status_code == 403:
                    print(colored(f"Access forbidden for URL: {url}", "red"))
                    blocked_urls.add(url)
                    return 0, ''

                elif response.status_code != 200:
                    print(colored(f"Failed to fetch URL: {url}, status code: {response.status_code}","red"))
                    blocked_urls.add(url)
                    return 0, ''

                parsed_url = urlparse(url)
                # print(colored(url,"blue"))

                if parsed_url.path.endswith('.pdf'):
                    pdf_reader = PdfReader(BytesIO(response.content))
                    num_pages = len(pdf_reader.pages)  
                    text = " ".join([pdf_reader.pages[i].extract_text() for i in range(num_pages)])

                elif 'youtube.com' in parsed_url.netloc or 'm.youtube.com' in parsed_url.netloc:
                    video_id = parsed_url.query.split('v=')[-1]
                    
                    try:
                        # Fetch video details using pytube
                        yt = YouTube(f"https://www.youtube.com/watch?v={video_id}")
                        title = yt.title
                        
                        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
                        auto_transcript = None
                        
                        for transcript in transcript_list:
                            if transcript.is_generated:
                                auto_transcript = transcript
                                break

                        if auto_transcript:
                            srt = auto_transcript.fetch()
                            text = " ".join([entry['text'] for entry in srt])
                            
                            # Prepend the title to the text
                            text = f"{title} - {text}"

                            detected_language = auto_transcript.language_code

                            if detected_language != 'en':
                                retry_count = 0
                                while retry_count < 5:
                                    try:
                                        response = openai.ChatCompletion.create(
                                            model="gpt-3.5-turbo-16k",
                                            temperature=0.3,
                                            messages=[
                                                {"role": "system", "content": f"""As a Language Expert. Directly, summarize text below from language: {detected_language} to English. 
                                                 Limit to not more than 3000 tokens
                                                Output format example:
                                                (url - title - summarized and detailed content
                                                 """},
                                                {"role": "user", "content":f""" text: {text}"""}
                                            ]
                                        )
                                        translated_text = response.choices[0].message['content'].strip()
                                        text = translated_text
                                        break  
                                    except Exception as trans_error:
                                        print(f"OpenAI Translation failed: {trans_error}")
                                        retry_count += 1

                                if retry_count == 5:
                                    text = f"This text is in '{detected_language}' - {text}"

                        else:
                            text = "No auto-generated captions available for this video."

                    except Exception as e:
                        print(f"An unexpected error occurred: {e}")

                elif parsed_url.path.endswith(('.csv', '.xlsx')):
                    # CSV and Excel handling
                    df = pd.read_csv(BytesIO(response.content)) if parsed_url.path.endswith('.csv') else pd.read_excel(BytesIO(response.content))
                    text = " ".join(df.to_string())

                elif parsed_url.path.endswith('.json'):
                    # JSON handling
                    text = json.dumps(response.json())

                elif parsed_url.path.endswith('.xml'):
                    # XML handling
                    soup = BeautifulSoup(response.content, 'xml')  # Note that the feature is set to 'xml'

                elif parsed_url.path.endswith('.pptx'):
                    # PPT handling
                    prs = Presentation(BytesIO(response.content))
                    text = " ".join([slide.notes_slide.notes_text_frame.text for slide in prs.slides])

                elif parsed_url.path.endswith('.docx'):
                    # Word handling
                    doc = Document(BytesIO(response.content))
                    text = " ".join([p.text for p in doc.paragraphs])

                elif parsed_url.path.endswith('.js'):
                    # JavaScript handling (simplified)
                    driver = webdriver.Firefox()
                    driver.get(url)
                    text = driver.page_source
                    driver.quit()

                else:
                    # HTML or other text-based content
                    soup = BeautifulSoup(response.text, 'html.parser')

                    # Remove script and style elements
                    for script_or_style in soup(["script", "style"]):
                        script_or_style.extract()
                    text = ' '.join(soup.stripped_strings)
                    text = ''.join(text)  # Convert to string

                # Normalize the text to 'NFKD' form
                normalized_text = unicodedata.normalize('NFKD', text)

                # Remove unwanted characters
                cleaned_text = re.sub(r'[^\w\s,.!?{}[]&@()% "\']', '', normalized_text)
                cleaned_text = re.sub(r'\.{2,}', '', cleaned_text)
                cleaned_text = re.sub(r'\s{2,}', ' ', cleaned_text)  # Remove more than one repetitive space
                # Count the number of tokens using tiktoken
                token_count_result_1 = self.token_count(cleaned_text)
                # print(f"token_count_result_1: {token_count_result_1}")
                # Initialize variables for truncation
                max_length_1 = 3000  # Replace with your desired max token count

                # Check if truncation is needed
                if token_count_result_1 > 0 and token_count_result_1 > max_length_1:
                    avg_token_length = len(cleaned_text) / token_count_result_1
                    approx_cut_index = int(avg_token_length * max_length_1)
                    truncated_text = cleaned_text[:approx_cut_index]

                    # Fine-tune to make sure we have exactly max_length tokens
                    prev_token_count = -1  # Initialize a variable to store the previous token count
                    while self.token_count(truncated_text) > max_length_1:
                        if prev_token_count == self.token_count(truncated_text):
                            break  # Exit the loop if the token count is not changing
                        prev_token_count = self.token_count(truncated_text)
                        truncated_text = truncated_text[:-1]  # Remove one character from the end

                    prev_token_count = -1  # Reset the variable for the next loop
                    while self.token_count(truncated_text) < max_length_1 and len(cleaned_text) > len(truncated_text):
                        if prev_token_count == self.token_count(truncated_text):
                            break  # Exit the loop if the token count is not changing
                        prev_token_count = self.token_count(truncated_text)
                        truncated_text += cleaned_text[len(truncated_text)]

                    cleaned_text = truncated_text

                # Count the number of tokens using tiktoken
                token_count_result2 = self.token_count(cleaned_text)
                print(colored(f"{url}, Token Count: {token_count_result2}","magenta"))
                # print(colored(f"{url}, cleaned_text: {cleaned_text}","magenta"))

                # Generate embeddings for the cleaned text
                cleaned_content = cleaned_text
                embedding = self.model([cleaned_content])[0].numpy()
                # print(colored(f"{url} Embedding","cyan"))

                return cleaned_content, embedding

            except Exception as e:
                print(colored(f"An error occurred: {e}","red"))
                if attempt < max_retries - 1:
                    print(colored(f"Retrying... ({attempt + 1})","red"))
                    time.sleep(delay)
                else:
                    print(colored("Max retries reached. Could not complete the operation.","red"))
                    return 0, ''

    # 3b 3 Use k-NN to sort internet links 
    def sorted_url_content_embedding_list(self, url_content_list, current_user_query_embedding):
        # Remove rows where the URL could not be scraped (i.e., cleaned_content is empty or None)
        filtered_url_content_list = [row for row in url_content_list if row[1]]

        # If the list is empty after filtering, return a message indicating that no links could be scraped
        if len(filtered_url_content_list) == 0:
            print(colored("No valid URLs found. All links could not be scraped.","red"))
            return []

        # Extract embeddings from the filtered list
        embeddings = [np.array(x[2]) for x in filtered_url_content_list]

        # Fit k-NN model
        knn = NearestNeighbors(n_neighbors=len(embeddings), algorithm='auto').fit(np.vstack(embeddings))
        distances, indices = knn.kneighbors([current_user_query_embedding])
        sorted_links = [(filtered_url_content_list[i][0], filtered_url_content_list[i][1]) for i in indices[0]]

        return sorted_links
    
#  4 COMBINING PIPELINES
    
    def combined_pipelines(self, raw_db_response, raw_internet_response, query):

        db_response = ""
        internet_response = ""

        # Check if the query contains a URL
        is_url_present = False
        internet_links = re.findall(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', query)
        if internet_links:
            is_url_present = True

        # Process database response only if no URL is present in the query
        if not is_url_present:
            db_response = self.text_cleaning(raw_db_response)

        # Process internet response regardless of the database response
        internet_response = self.text_cleaning(raw_internet_response)

        # Combine the two responses
        combined_response = internet_response + db_response

        return db_response, internet_response, combined_response

    def token_count(self, string: str, encoding_name: str = "cl100k_base") -> int:
        encoding = tiktoken.get_encoding(encoding_name)
        num_tokens = len(encoding.encode(string, disallowed_special=(encoding.special_tokens_set - {''})))
        return num_tokens
    

# TEXT CLEANING
# ================================================================================================ #

    def text_cleaning(self, text):

        # print(text)
        text=str(text)
        # Decode Unicode escape sequences
        text = bytes(text, "utf-8").decode("unicode_escape")
        
        # Extract URLs
        urls = re.findall(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', text)

        # Replace URLs with placeholders
        for i, url in enumerate(urls):
            placeholder = f"URL_PLACEHOLDER_{i}"
            text = text.replace(url, placeholder)
        
        # Cleaning steps
        text = ''.join(char for char in text if ord(char) < 128)
        text = re.sub(r'[^a-zA-Z0-9,.!?{}[]&@()% "\'_]', '', text)
        text = re.sub(r'(None , ){2,}', 'None , ', text)
        text = re.sub(r'(, ){2,}', ', ', text)

        # Count the number of tokens using your token_count function
        token_count_result_2 = self.token_count(text)

        # Initialize variables for truncation
        max_length = 5000  # Replace with your desired max token count

        # Check if truncation is needed
        if token_count_result_2 > 0 and token_count_result_2 > max_length:
            avg_token_length = len(text) / token_count_result_2
            approx_cut_index = int(avg_token_length * max_length)
            truncated_text = text[:approx_cut_index]

            # Fine-tune to make sure we have exactly max_length tokens
            prev_token_count = -1  # Initialize a variable to store the previous token count
            while self.token_count(truncated_text) > max_length:
                if prev_token_count == self.token_count(truncated_text):
                    break  # Exit the loop if the token count is not changing
                prev_token_count = self.token_count(truncated_text)
                truncated_text = truncated_text[:-1]  # Remove one character from the end

            prev_token_count = -1  # Reset the variable for the next loop
            while self.token_count(truncated_text) < max_length and len(text) > len(truncated_text):
                if prev_token_count == self.token_count(truncated_text):
                    break  # Exit the loop if the token count is not changing
                prev_token_count = self.token_count(truncated_text)
                truncated_text += text[len(truncated_text)]  # Add one more character from the original text

            text = truncated_text  # Update text with the truncated version

        # Put URLs back in
        for i, url in enumerate(urls):
            placeholder = f"URL_PLACEHOLDER_{i}"
            text = text.replace(placeholder, url)
            
        return text
         
# 5 COMBINED RESPONSE  
    
    def generate_final_response(self, db_response, internet_response, query):

        # # Generate the current timestamp
        # self.current_timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-16k", # gpt-4
            temperature=0.7,
            messages=[
                {"role": "system",
                "content": f"""
                    You have memory and real-time internet access. 
                    Your goal is to synthesize the information from either or both sources to provide the most accurate and relevant response to the current query. Just ignore if the output is '-' or empty. 
                    Provide the link/s and title at the end for reference. 
                    Provide it in detail as possible while limiting to 6000 tokens of output.
                """
                },
                {
                "role": "user",
                "content": f"""
                    (Memory-based Output: {db_response}.)
                    (Internet-based Output (Latest/Real-time based on curret query): {internet_response}.)
                    (Current Query: {query}.)
                    (Current Timestamp: {self.current_timestamp}.)
                """
                }
                ]
            )
        final_response = response.choices[0].message['content'].strip()
        final_response = str(final_response)
        
        return final_response

# SUMMARIZER

    def generate_conversation_summary(self, final_response, query):

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-16k", # gpt-4
            temperature=0.7,
            messages=[
                {"role": "system",
                "content": f"""
                At like we are in a converstation. Your role is to summarize the current conversation between Hafiz (user) and Skypea (yourself) in a way that captures the essence of the discussion, the questions asked, and the information provided.
                This summary will be used for future reference.
                Limited to 1000 tokens of output.
                """
                },
                {
                "role": "user",
                "content": f"""
                    (My Current Query: {query}.)
                    (Your Reply: {final_response}.)
                    
                    Example Summary Format:
                    "You : query.
                    Me : your reply.
                    Reference: 
                    Link_1 - Title_1
                    Link_2 - Title_2 "
                """
                }
                ]
            )
        summary = response.choices[0].message['content'].strip()
        
        return summary

# STORING DATA

    def store_data_in_db(self, query, current_user_query_embedding, db_response, internet_response, final_response, summary):
        try:
            # Generate embeddings for db_response, internet_response, and combined_response
            db_response_embedding = self.model([db_response])[0].numpy().tolist()
            internet_response_embedding = self.model([internet_response])[0].numpy().tolist()
            final_response_embedding = self.model([final_response])[0].numpy().tolist()

            # Convert numpy arrays to lists for database storage (if they are numpy arrays)
            current_user_query_embedding_list = current_user_query_embedding.tolist() if isinstance(current_user_query_embedding, np.ndarray) else current_user_query_embedding

            # Inserting the data and embeddings into the database
            self.cur.execute("""
                INSERT INTO skypea_db (
                    user_query,
                    db_response,
                    internet_response,
                    combined_response,
                    user_query_embedding,
                    db_response_embedding,
                    internet_response_embedding,
                    combined_response_embedding,
                    summary
                ) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s) RETURNING id, timestamp;
            """, (query, db_response, internet_response, final_response, current_user_query_embedding_list, db_response_embedding, internet_response_embedding, final_response_embedding, summary))

            # Committing the changes to the database
            self.conn.commit()

            # Fetching the returned id and timestamp
            iid, timestamp = self.cur.fetchone()

            # Return the ID and timestamp for any further use
            return iid, timestamp, current_user_query_embedding

        except psycopg2.Error as e:
            print(colored("Error: Unable to store data in the database", "red"))
            print(colored({e},"red"))
            # Rollback in case of any error
            self.conn.rollback()
            return None, None, None

# MAIN LOOP

if __name__ == "__main__":
    gpt_instance = gpt()

    print(colored("Welcome to Skypea! What's your query?", "green"))
    print(colored("Enter 'Exit' = Quit, Enter 'Done' = Query","yellow"))

    while True:
        
        query = ''
        print(colored("Please enter your query (type '?' on a new line when finished):","yellow"))
        while True:
            line = input()
            if line.strip().lower() == '?':
                break
            elif line.strip().lower() == 'exit':
                gpt.close_connection()
                exit()
            query += line + '\n'

        clear_terminal()

        main_start_time = time.time()  # Record the start time

        # Generate Current Query Embeddings
        section_start_time = time.time()
        current_user_query_embedding = gpt_instance.generate_current_query_embedding(query)

        # Initialize k-NN Model
        section_start_time = time.time()
        gpt_instance.initialize_knn_model()
        
        pipeline_start_time = time.time()

        # Submitting the Database Pipeline function to the ThreadPoolExecutor
        memory, internet_query = gpt_instance.fetch_and_clean_from_db(current_user_query_embedding, query)
        raw_db_response = memory
        # print(raw_db_response)
        db_end_time = time.time()
        print(colored(f"3a DATABASE ROUTE: {db_end_time - pipeline_start_time:.2f} seconds", "yellow"))

        # Wait for the Internet Pipeline to complete and get the result
        sorted_links = gpt_instance.internet_route(current_user_query_embedding, query, internet_query)
        raw_internet_response = sorted_links
        # print(raw_internet_response)
        internet_end_time = time.time()
        print(colored(f"3b INTERNET ROUTE: {internet_end_time - pipeline_start_time:.2f} seconds", "yellow"))

        # Combined Pipelines
        section_start_time = time.time()
        db_response, internet_response, combined_response = gpt_instance.combined_pipelines(raw_db_response, raw_internet_response, query)
        
        # Generate Response
        final_response = gpt_instance.generate_final_response(db_response, internet_response, query)
        print(colored(final_response, "green"))
        print(colored(f"5 FINAL RESPONSE: {time.time() - section_start_time:.2f} seconds", "yellow"))

        # Generate Conversation Summary
        summary = gpt_instance.generate_conversation_summary(final_response, query)
        print(colored(f"6 SUMMARY: {time.time() - section_start_time:.2f} seconds", "yellow"))

        # STORING DATA
        gpt_instance.store_data_in_db(query, current_user_query_embedding, db_response, internet_response, final_response, summary)
        
        main_end_time = time.time()  # Record the end time
        print(colored(f"Total time taken for main loop: {main_end_time - main_start_time:.2f} seconds", "yellow"))